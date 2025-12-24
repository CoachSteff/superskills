"""
Presenter.py - Slide deck generation from markdown using python-pptx.
"""
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Literal, Optional

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not available - install with: pip install python-pptx")

try:
    import markdown
    MARKDOWN_AVAILABLE = True
except ImportError:
    MARKDOWN_AVAILABLE = False
    print("Warning: markdown not available - install with: pip install markdown")


@dataclass
class PresentationResult:
    """Result from a presentation operation."""
    source_file: str
    output_file: str
    slide_count: int
    format: str
    file_size_mb: float
    processing_time_seconds: float
    timestamp: str = None

    def __post_init__(self):
        if not self.timestamp:
            self.timestamp = datetime.now().isoformat()


class Presenter:
    """Slide deck generation from markdown using python-pptx."""

    # Presentation themes
    THEMES = {
        "default": {
            "background": RGBColor(255, 255, 255),
            "title_color": RGBColor(31, 78, 121),
            "text_color": RGBColor(0, 0, 0)
        },
        "dark": {
            "background": RGBColor(30, 30, 30),
            "title_color": RGBColor(255, 200, 87),
            "text_color": RGBColor(255, 255, 255)
        },
        "professional": {
            "background": RGBColor(245, 245, 245),
            "title_color": RGBColor(0, 51, 102),
            "text_color": RGBColor(51, 51, 51)
        }
    }

    def __init__(
        self,
        output_dir: str = "output/presentations",
        theme: str = "default",
        verbose: bool = True
    ):
        """Initialize Presenter.

        Args:
            output_dir: Directory to save presentations
            theme: Presentation theme (default, dark, professional)
            verbose: Enable verbose logging
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.verbose = verbose

        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is required. Install with: pip install python-pptx")

        if not MARKDOWN_AVAILABLE:
            raise ImportError("markdown is required. Install with: pip install markdown")

        if theme not in self.THEMES:
            raise ValueError(f"Unknown theme: {theme}. Choose from: {list(self.THEMES.keys())}")

        self.theme = self.THEMES[theme]
        self.prs = None

    def create_from_markdown(
        self,
        markdown_file: str,
        output_name: Optional[str] = None,
        title_slide: bool = True
    ) -> PresentationResult:
        """Create presentation from markdown file.

        Args:
            markdown_file: Path to markdown file
            output_name: Output filename (without extension)
            title_slide: Include title slide

        Returns:
            PresentationResult with presentation details
        """
        markdown_path = Path(markdown_file)
        if not markdown_path.exists():
            raise FileNotFoundError(f"Markdown file not found: {markdown_file}")

        if self.verbose:
            print(f"Creating presentation from: {markdown_path.name}")

        start_time = datetime.now()

        # Read markdown content
        with open(markdown_path, 'r', encoding='utf-8') as f:
            content = f.read()

        # Parse markdown into slides
        slides_data = self._parse_markdown(content)

        # Create presentation
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

        # Add title slide if requested
        if title_slide and slides_data:
            first_slide = slides_data[0]
            if first_slide["type"] == "title":
                self._add_title_slide(first_slide["title"], first_slide.get("subtitle", ""))
                slides_data = slides_data[1:]

        # Add content slides
        for slide_data in slides_data:
            if slide_data["type"] == "title":
                self.add_slide(
                    title=slide_data["title"],
                    content=slide_data.get("subtitle", ""),
                    layout="title"
                )
            elif slide_data["type"] == "content":
                self.add_slide(
                    title=slide_data["title"],
                    content=slide_data["content"],
                    layout="bullet"
                )

        # Save presentation
        output_file = self.output_dir / (output_name or f"{markdown_path.stem}.pptx")
        self.prs.save(str(output_file))

        processing_time = (datetime.now() - start_time).total_seconds()
        file_size_mb = output_file.stat().st_size / (1024 * 1024)

        result = PresentationResult(
            source_file=str(markdown_path),
            output_file=str(output_file),
            slide_count=len(self.prs.slides),
            format="pptx",
            file_size_mb=file_size_mb,
            processing_time_seconds=processing_time
        )

        if self.verbose:
            print(f"✓ Created presentation: {len(self.prs.slides)} slides")
            print(f"✓ Saved to: {output_file}")

        return result

    def add_slide(
        self,
        title: str,
        content: str = "",
        layout: Literal["title", "bullet", "two_column"] = "bullet",
        image_path: Optional[str] = None
    ) -> int:
        """Add a slide to the presentation.

        Args:
            title: Slide title
            content: Slide content (bullet points or text)
            layout: Slide layout type
            image_path: Path to image to add (optional)

        Returns:
            Slide index
        """
        if not self.prs:
            self.prs = Presentation()

        if layout == "title":
            slide = self._add_title_slide(title, content)
        elif layout == "bullet":
            slide = self._add_bullet_slide(title, content)
        elif layout == "two_column":
            slide = self._add_two_column_slide(title, content)
        else:
            raise ValueError(f"Unknown layout: {layout}")

        # Add image if provided
        if image_path and Path(image_path).exists():
            self._add_image_to_slide(slide, image_path)

        return len(self.prs.slides) - 1

    def export_to_pdf(
        self,
        pptx_file: str,
        output_name: Optional[str] = None
    ) -> str:
        """Export presentation to PDF.

        Note: This requires LibreOffice or PowerPoint to be installed.

        Args:
            pptx_file: Path to PowerPoint file
            output_name: Output PDF filename

        Returns:
            Path to PDF file
        """
        pptx_path = Path(pptx_file)
        if not pptx_path.exists():
            raise FileNotFoundError(f"PowerPoint file not found: {pptx_file}")

        if self.verbose:
            print(f"Exporting to PDF: {pptx_path.name}")

        output_file = self.output_dir / (output_name or f"{pptx_path.stem}.pdf")

        # Try using LibreOffice for conversion
        try:
            import subprocess
            subprocess.run(
                [
                    "soffice",
                    "--headless",
                    "--convert-to", "pdf",
                    "--outdir", str(self.output_dir),
                    str(pptx_path)
                ],
                check=True,
                capture_output=True
            )

            if self.verbose:
                print(f"✓ Exported to PDF: {output_file}")

            return str(output_file)

        except (subprocess.CalledProcessError, FileNotFoundError):
            raise RuntimeError(
                "PDF export requires LibreOffice. Install with: brew install libreoffice"
            )

    def _parse_markdown(self, content: str) -> List[Dict]:
        """Parse markdown content into slide data."""
        slides = []
        current_slide = None

        lines = content.split('\n')

        for line in lines:
            # H1 = New title slide
            if line.startswith('# '):
                if current_slide:
                    slides.append(current_slide)
                current_slide = {
                    "type": "title",
                    "title": line[2:].strip(),
                    "subtitle": ""
                }

            # H2 = New content slide
            elif line.startswith('## '):
                if current_slide:
                    slides.append(current_slide)
                current_slide = {
                    "type": "content",
                    "title": line[3:].strip(),
                    "content": []
                }

            # Bullet points
            elif line.strip().startswith('- ') or line.strip().startswith('* '):
                if current_slide and current_slide["type"] == "content":
                    current_slide["content"].append(line.strip()[2:])

            # Regular text
            elif line.strip() and current_slide:
                if current_slide["type"] == "title":
                    current_slide["subtitle"] += line.strip() + " "
                elif current_slide["type"] == "content":
                    current_slide["content"].append(line.strip())

        if current_slide:
            slides.append(current_slide)

        return slides

    def _add_title_slide(self, title: str, subtitle: str = "") -> object:
        """Add title slide."""
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)

        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme["background"]

        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.color.rgb = self.theme["title_color"]
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True

        # Subtitle
        if subtitle and len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_frame = subtitle_shape.text_frame
            subtitle_frame.paragraphs[0].font.color.rgb = self.theme["text_color"]
            subtitle_frame.paragraphs[0].font.size = Pt(24)

        return slide

    def _add_bullet_slide(self, title: str, content: str) -> object:
        """Add bullet point slide."""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)

        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme["background"]

        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.color.rgb = self.theme["title_color"]
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True

        # Content
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()

        if isinstance(content, list):
            for i, bullet in enumerate(content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.color.rgb = self.theme["text_color"]
                p.font.size = Pt(18)
        else:
            p = text_frame.paragraphs[0]
            p.text = content
            p.font.color.rgb = self.theme["text_color"]
            p.font.size = Pt(18)

        return slide

    def _add_two_column_slide(self, title: str, content: str) -> object:
        """Add two-column slide."""
        # Use blank layout and add shapes manually
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme["background"]

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.color.rgb = self.theme["title_color"]
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True

        return slide

    def _add_image_to_slide(self, slide: object, image_path: str):
        """Add image to slide."""
        left = Inches(6)
        top = Inches(2)
        height = Inches(4)

        slide.shapes.add_picture(
            image_path,
            left, top,
            height=height
        )


def create_presentation_from_markdown(
    markdown_file: str,
    output_dir: str = "output/presentations",
    **kwargs
) -> PresentationResult:
    """Convenience function to create presentation from markdown.

    Args:
        markdown_file: Path to markdown file
        output_dir: Output directory
        **kwargs: Additional arguments for Presenter

    Returns:
        PresentationResult
    """
    presenter = Presenter(output_dir=output_dir)
    return presenter.create_from_markdown(markdown_file, **kwargs)
